import os
import json
import requests
import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PPTGenerator:
    def __init__(self, api_key):
        self.api_key = api_key
        if not self.api_key:
            raise ValueError("API key is required")
        genai.configure(api_key=self.api_key)
        # Corrected model names to current valid models
        self.model = genai.GenerativeModel("gemini-1.5-flash")
        self.model_vision = genai.GenerativeModel("gemini-1.5-pro")
        self.presentation = Presentation()

    def generate_content_outline(self, topic, num_slides=5):
        """Generate content outline using the Gemini model."""
        prompt = f"""Create a content outline for a presentation on the topic '{topic}' with {num_slides} slides.
        Provide a title slide, several content slides, and a conclusion slide.
        Return the response as a valid JSON array with the following structure for each slide:
        [
            {{
                "title": "Slide Title",
                "content": "Main content points as bullet points, separated by newlines.",
                "slide_type": "title|content|image|conclusion"
            }}
        ]
        The response must be only the JSON array and nothing else.
        """
        try:
            response = self.model.generate_content(prompt)
            content = response.text.strip()
            if "```json" in content:
                content = content.split("```json")[1].split("```")[0].strip()
            elif "```" in content:
                content = content.split("```")[1].strip()

            # Fixed typo: startswith, not startwith
            if not content.startswith('[') or not content.endswith(']'):
                print("Error: Model did not return a valid JSON array string.")
                return None
            
            return json.loads(content)
        except json.JSONDecodeError as e:
            print(f"JSON Decode Error: {e}\nRaw content: {content}")
            return None
        except Exception as e:
            print(f"Error generating content outline: {e}")
            return None

    def generate_image_description(self, slide_content):
        """Generates a search query for an image based on slide content."""
        prompt = f"""
        Based on the following slide content, suggest a relevant image search query.
        Content: {slide_content}
        Return only a brief descriptive phrase suitable for an image search (max 5 words).
        """
        try:
            response = self.model_vision.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            print(f"Error generating image description: {e}")
            return "abstract technology background" # Fallback query

    def download_image(self, query, save_path="temp_img.jpg"):
        """Downloads an image from Pexels."""
        try:
            pexels_api_key = os.getenv("PEXELS_API_KEY")
            if not pexels_api_key:
                print("PEXELS_API_KEY not found in environment variables.")
                return None

            url = "[https://api.pexels.com/v1/search](https://api.pexels.com/v1/search)"
            headers = {"Authorization": pexels_api_key}
            params = {"query": query, "per_page": 1, "orientation": "landscape"}
            
            response = requests.get(url, headers=headers, params=params)
            response.raise_for_status()
            data = response.json()

            # Fixed key: "photos" not "photo"
            if not data.get("photos"):
                print(f"No images found for query: '{query}'")
                return None
            
            image_url = data["photos"][0]["src"]["original"]
            image_response = requests.get(image_url)
            image_response.raise_for_status()

            with open(save_path, 'wb') as f:
                f.write(image_response.content)
            return save_path
        except Exception as e:
            print(f"Error downloading image: {e}")
            return None

    def create_title_slide(self, title, subtitle):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Fixed typo: title_shape, not title.shape
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].font.italic = True
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)
            subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def create_content_slide(self, title, content, include_image=False):
        slide_layout = self.presentation.slide_layouts[1] # Title and Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear() # Clear existing text
        text_frame.word_wrap = True

        for point in content.split('\n'):
            p = text_frame.add_paragraph()
            p.text = point.strip('- ').strip()
            p.font.size = Pt(18)
            p.level = 0
        
        if include_image:
            image_desc = self.generate_image_description(content)
            img_path = self.download_image(image_desc)
            if img_path:
                # Adjust content box to make space for the image
                content_shape.width = Inches(5.5)
                left = Inches(4.25)
                top = Inches(1.75)
                height = Inches(4.0)
                slide.shapes.add_picture(img_path, left, top, height=height)
                os.remove(img_path)

    # Corrected method signature to accept content
    def create_image_slide(self, title, content, img_query):
        slide_layout = self.presentation.slide_layouts[5] # Title Only layout
        slide = self.presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Fixed typo: title_shape, not title.shape
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if content:
             # Add a text box for the content/caption
            left, top, width, height = Inches(0.5), Inches(6), Inches(9), Inches(1)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = content
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

        img_path = self.download_image(img_query)
        if img_path:
            # Center the image on the slide
            left = Inches(1)
            top = Inches(1.5)
            height = Inches(4.5)
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
            os.remove(img_path)

    def generate_presentation(self, topic, num_slides=5, output_file="presentation.pptx"):
        print("Generating content outline...")
        outline = self.generate_content_outline(topic, num_slides)
        if not outline or len(outline) < 1:
            print("Failed to generate a valid content outline.")
            return

        print("Creating title slide...")
        self.create_title_slide(outline[0]['title'], outline[0].get('content', ''))

        # Corrected loop logic using enumerate
        for i, slide_info in enumerate(outline[1:], start=1):
            slide_type = slide_info.get('slide_type', 'content').lower()
            title = slide_info.get('title', '')
            content = slide_info.get('content', '')

            print(f"Creating slide {i+1}/{len(outline)}: '{title}' (Type: {slide_type})")

            if slide_type == 'title':
                self.create_title_slide(title, content)
            elif slide_type == 'image':
                img_query = self.generate_image_description(content)
                # Corrected method call to pass all required arguments
                self.create_image_slide(title, content, img_query)
            elif slide_type == 'conclusion':
                self.create_content_slide(title, content, include_image=False)
            else: # Default to 'content'
                self.create_content_slide(title, content, include_image=True)

        self.presentation.save(output_file)
        print(f"\nPresentation saved as {output_file}")
        return output_file

if __name__ == "__main__":
    load_dotenv()
    print("Generation Initiated...")

    try:
        # Ensure you have GOOGLE_API_KEY and PEXELS_API_KEY in your .env file
        generator = PPTGenerator(os.getenv("GOOGLE_API_KEY"))
        topic = "Animal Kingdom"
        output_file = generator.generate_presentation(topic, num_slides=7, output_file="AI_Presentation.pptx")
        if output_file:
            print(f"Presentation generated successfully: {output_file}")
    except Exception as e:
        print(f"An error occurred: {e}")